import os
import json
import logging
import numpy as np
import pytesseract
import textract
import tiktoken
import nltk
import base64
import asyncio
import websockets
import requests
import httpx 
import tempfile
from io import BytesIO
from fastapi import FastAPI, WebSocket, UploadFile, Request, File
from fastapi.responses import HTMLResponse, JSONResponse
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from pdf2image import convert_from_bytes
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from fastapi.middleware.cors import CORSMiddleware
from starlette.websockets import WebSocketDisconnect
from typing import List, Generator
from pydantic import Field
from contextlib import asynccontextmanager
from nltk.tokenize import sent_tokenize
from dotenv import load_dotenv
from striprtf.striprtf import rtf_to_text
from groq import Groq
from pydub import AudioSegment
from pydantic import BaseModel, ValidationError
from fastapi.staticfiles import StaticFiles


# Importing from llama_index.core
from llama_index.core import VectorStoreIndex, Document, StorageContext
from llama_index.core.embeddings import BaseEmbedding

# Ensure NLTK data is downloaded
nltk.download('punkt')

# Load environment variables
load_dotenv()

# Configuration for logging
logging.basicConfig(
    level=logging.DEBUG,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("app.log"),
        logging.StreamHandler()
    ]
)

# Set up API keys
NVIDIA_API_KEY_EMBEDQA = os.getenv("NVIDIA_API_KEY_EMBEDQA")
NVIDIA_API_KEY_LLAMA = os.getenv("NVIDIA_API_KEY_LLAMA")
GROQ_API_KEY = os.getenv("GROQ_API_KEY")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")

# Check for missing API keys
if not NVIDIA_API_KEY_EMBEDQA:
    logging.error("NVIDIA_API_KEY_EMBEDQA is not set.")
else:
    logging.info("NVIDIA_API_KEY_EMBEDQA loaded successfully.")

if not NVIDIA_API_KEY_LLAMA:
    logging.error("NVIDIA_API_KEY_LLAMA is not set.")
else:
    logging.info("NVIDIA_API_KEY_LLAMA loaded successfully.")

if not GROQ_API_KEY:
    logging.error("GROQ_API_KEY is not set.")
else:
    logging.info("GROQ_API_KEY loaded successfully.")

if not OPENAI_API_KEY:
    logging.error("OPENAI_API_KEY is not set.")
else:
    logging.info("OPENAI_API_KEY loaded successfully.")

# URLs for NVIDIA APIs
LLAMA_URL = "https://integrate.api.nvidia.com/v1"
EMBEDQA_URL = "https://integrate.api.nvidia.com/v1"
REALTIME_API_URL = "wss://api.openai.com/v1/realtime?model=gpt-4o-realtime-preview-2024-10-01"

# Paths for file directories
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
DATABASE_DIR = os.path.join(BASE_DIR, "database")
STATIC_DIR = os.path.join(BASE_DIR, "static")
TEMPLATES_DIR = os.path.join(BASE_DIR, "templates")

# Ensure directories exist
os.makedirs(DATABASE_DIR, exist_ok=True)
os.makedirs(STATIC_DIR, exist_ok=True)
os.makedirs(TEMPLATES_DIR, exist_ok=True)

# Minimum text length to consider extracted text as valid
MIN_TEXT_LENGTH = 50

# ------------------------ Embedding Generation Functions ------------------------

def get_text_embedding(text: str) -> List[float]:
    """Generates an embedding for the given text using NVIDIA's Embedding API."""
    try:
        # NVIDIA API endpoint for embeddings
        url = f"{EMBEDQA_URL}/embeddings"

        # Set up the headers with the API key
        headers = {
            "Authorization": f"Bearer {NVIDIA_API_KEY_EMBEDQA}",
            "Content-Type": "application/json",
        }

        # Prepare the payload
        payload = {
            "input": [text],
            "model": "nvidia/nv-embedqa-mistral-7b-v2",
            "encoding_format": "float",
            "input_type": "query",
            "truncate": "NONE"
        }

        # Make the POST request
        response = requests.post(url, headers=headers, json=payload)

        # Check for successful response
        if response.status_code != 200:
            logging.error(f"Failed to get embedding: {response.status_code} {response.text}")
            return []

        response_json = response.json()
        embedding = response_json["data"][0]["embedding"]
        logging.info(f"Generated embedding for text: {text[:50]}... with size {len(embedding)}")
        return embedding
    except Exception as e:
        logging.error(f"Unexpected error during embedding generation: {e}")
        return []

# ------------------------ Custom Embedding Class for llama_index ------------------------

class NVIDIACustomEmbedding(BaseEmbedding):
    """Custom embedding class for llama_index using NVIDIA's Embedding API."""

    def __init__(self):
        super().__init__()

    def _get_text_embedding(self, text: str) -> List[float]:
        return get_text_embedding(text)

    def _get_query_embedding(self, text: str) -> List[float]:
        return get_text_embedding(text)

    async def _aget_query_embedding(self, text: str) -> List[float]:
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(None, self._get_query_embedding, text)

    def embed_documents(self, texts: List[str]) -> List[List[float]]:
        embeddings = []
        for text in texts:
            try:
                embedding = self._get_text_embedding(text)
                embeddings.append(embedding)
            except Exception as e:
                logging.error(f"Error generating embedding for document: {text[:50]}... Error: {e}")
                embeddings.append([])
        return embeddings

    def embed_query(self, text: str) -> List[float]:
        try:
            embedding = self._get_query_embedding(text)
            return embedding
        except Exception as e:
            logging.error(f"Error generating embedding for query: {text[:50]}... Error: {e}")
            return []

# ------------------------ Initialize the Vector Store Index ------------------------

vector_index = None
embedding_model = None  # Added this line to initialize embedding_model

def split_text_into_token_limited_chunks(text: str, max_tokens: int = 300) -> List[str]:
    """Splits text into chunks limited by the specified maximum number of tokens."""
    try:
        tokenizer = tiktoken.get_encoding("cl100k_base")
    except Exception as e:
        logging.error(f"Error initializing tokenizer: {e}")
        return []

    sentences = sent_tokenize(text)
    chunks = []
    current_chunk = ""
    current_tokens = 0

    for sentence in sentences:
        sentence_tokens = len(tokenizer.encode(sentence))
        if sentence_tokens > max_tokens:
            words = sentence.split()
            sub_sentence = ""
            for word in words:
                word_tokens = len(tokenizer.encode(word))
                if current_tokens + word_tokens > max_tokens:
                    if sub_sentence:
                        chunks.append(sub_sentence.strip())
                        logging.debug(f"Sub-sentence appended with {current_tokens} tokens.")
                    sub_sentence = word
                    current_tokens = word_tokens
                else:
                    sub_sentence += " " + word
                    current_tokens += word_tokens
            if sub_sentence:
                chunks.append(sub_sentence.strip())
                logging.debug(f"Sub-sentence appended with {current_tokens} tokens.")
                current_chunk = ""
                current_tokens = 0
        elif current_tokens + sentence_tokens > max_tokens:
            if current_chunk:
                chunks.append(current_chunk.strip())
                logging.debug(f"Chunk appended with {current_tokens} tokens.")
            current_chunk = sentence
            current_tokens = sentence_tokens
        else:
            current_chunk += " " + sentence
            current_tokens += sentence_tokens

    if current_chunk:
        chunks.append(current_chunk.strip())
        logging.debug(f"Final chunk appended with {current_tokens} tokens.")

    logging.info(f"Text successfully split into {len(chunks)} chunks.")
    return chunks

def index_document(text: str, doc_id: str):
    """Indexes a document by splitting it into chunks and processing each chunk."""
    chunks = split_text_into_token_limited_chunks(text, max_tokens=300)
    index_document_chunks(chunks, doc_id)

def index_document_chunks(chunks: List[str], doc_id: str):
    """Indexes individual document chunks."""
    global vector_index, embedding_model  # Added 'embedding_model' here

    for idx, chunk in enumerate(chunks):
        chunk_doc_id = f"{doc_id}_chunk_{idx}"
        try:
            tokenizer = tiktoken.get_encoding("cl100k_base")
            token_count = len(tokenizer.encode(chunk))
        except Exception as e:
            logging.error(f"Error initializing tokenizer for chunk {chunk_doc_id}: {e}")
            continue

        logging.debug(f"Processing chunk {chunk_doc_id} with {token_count} tokens.")

        if token_count > 512:
            logging.error(f"Chunk {chunk_doc_id} exceeds 512 tokens. Skipping.")
            continue

        try:
            embedding = embedding_model.embed_query(chunk)

            if not embedding:
                logging.error(f"No embedding generated for chunk {chunk_doc_id}. Skipping.")
                continue

            doc = Document(text=chunk, doc_id=chunk_doc_id, embedding=embedding)
            vector_index.insert(doc)
            logging.info(f"Document chunk {chunk_doc_id} inserted into the index.")
        except AttributeError as ae:
            logging.error(f"AttributeError: {ae}. Ensure that 'embedding_model' is correctly initialized.")
            continue
        except Exception as e:
            logging.error(f"Failed to insert chunk {chunk_doc_id} into the index: {e}")
            continue

def initialize_vector_store_index():
    """Initializes the VectorStoreIndex with the custom embedding model and indexes existing documents."""
    global vector_index, embedding_model  # Make sure embedding_model is global

    try:
        # Initialize the custom embedding model only if it hasn’t been initialized
        if embedding_model is None:
            embedding_model = NVIDIACustomEmbedding()

        # Ensure embedding_model is available before proceeding
        if not embedding_model:
            logging.error("Embedding model initialization failed.")
            return

        # Initialize storage context
        storage_context = StorageContext.from_defaults()

        # Initialize VectorStoreIndex with the embedding model
        vector_index = VectorStoreIndex(
            nodes=[],  # Start with an empty list of nodes
            storage_context=storage_context,
            embed_model=embedding_model  # Pass initialized embedding model
        )
        logging.info("Initialized VectorStoreIndex with custom embedding model.")

        # Load and index documents from the database directory
        for filename in os.listdir(DATABASE_DIR):
            if filename.endswith('.txt'):
                file_path = os.path.join(DATABASE_DIR, filename)
                try:
                    with open(file_path, 'r', encoding='utf-8') as f:
                        text = f.read()
                        # Index the document using chunks to avoid issues with large text blocks
                        index_document(text, doc_id=os.path.splitext(filename)[0])
                except Exception as e:
                    logging.error(f"Failed to read document {filename}: {e}")

        logging.info("VectorStoreIndex initialized with documents.")

    except Exception as e:
        logging.error(f"Failed to initialize VectorStoreIndex: {e}")
        vector_index = None  # Ensure vector_index remains None if initialization fails

# ------------------------ Application Lifespan Event ------------------------

@asynccontextmanager
async def lifespan(app: FastAPI):
    """Handles the application's lifespan events."""
    global vector_index
    initialize_vector_store_index()
    if vector_index is None:
        logging.error("VectorStoreIndex was not initialized properly.")
    else:
        logging.info("VectorStoreIndex initialized successfully.")
    yield
    pass

app = FastAPI(lifespan=lifespan)

# Add CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Mount static files and setup templates
app.mount("/static", StaticFiles(directory=STATIC_DIR), name="static")
templates = Jinja2Templates(directory=TEMPLATES_DIR)

# ------------------------ Helper Functions ------------------------

def is_text_valid(text: str) -> bool:
    """Checks if the extracted text is valid."""
    if not text.strip():
        return False
    if len(text.strip()) < MIN_TEXT_LENGTH:
        return False
    # Check ratio of non-alphanumeric characters
    alnum_count = sum(c.isalnum() for c in text)
    total_count = len(text)
    if total_count == 0 or (alnum_count / total_count) < 0.1:
        return False
    return True

def save_processed_document_to_db(doc_id: str, text: str):
    """Save the processed document text to DATABASE_DIR with tool usage information."""
    db_file_path = os.path.join(DATABASE_DIR, f"{doc_id}.txt")

    try:
        with open(db_file_path, "w", encoding="utf-8") as f:
            f.write(text)
        logging.info(f"Document {doc_id} saved to the index.")
    except Exception as e:
        logging.error(f"Failed to save document {doc_id} to database: {e}")

# ------------------------ Image Processing Functions ------------------------

def extract_data_with_groq(image_content, groq_api_key):
    """
    Extracts all visible text from an image using Groq's llama-3.2-11b-vision-preview model.
    """
    try:
        # Encode image content to base64
        base64_image = base64.b64encode(image_content).decode('utf-8')

        client = Groq(api_key=groq_api_key)  # Pass the API key to the Groq client

        # Define the new prompt
        prompt_text = (
            "Extract and list all visible text from the image, keeping any apparent structure, labels, "
            "numbers, and annotations intact.\n"
            "Transcribe every visible text element, label, and number exactly as it appears.\n"
            "Capture all text elements and numbers in the order they appear visually, including any names, "
            "categories, population figures, and other numerical data.\n"
            "List each entry on a new line to reflect the visual structure, and retain any relationships "
            "and groupings implied by the layout."
        )

        chat_completion = client.chat.completions.create(
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt_text},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{base64_image}",
                            },
                        },
                    ],
                }
            ],
            model="llama-3.2-11b-vision-preview",
        )

        extracted_text = chat_completion.choices[0].message.content

        return extracted_text

    except Exception as e:
        logging.error(f"Error during Groq llama-3.2-11b-vision-preview extraction: {e}")
        return ""

async def process_image_content(image_content: bytes, source_info: str) -> str:
    """
    Processes an image by performing pytesseract and Groq's llama-3.2-11b-vision-preview extraction.
    Returns a string containing grouped information from both methods.
    """
    try:
        # Load image
        with Image.open(BytesIO(image_content)) as img:
            # Adjust DPI metadata without resizing
            img.info['dpi'] = (300, 300)
            img_byte_arr = BytesIO()
            img.save(img_byte_arr, format='PNG', dpi=(300, 300))
            processed_image_bytes = img_byte_arr.getvalue()

        # Initialize extracted data list
        extracted_data = []

        # OCR extraction using pytesseract
        ocr_text = pytesseract.image_to_string(Image.open(BytesIO(processed_image_bytes)), lang='deu+eng')
        if is_text_valid(ocr_text):
            extracted_data.append(f"pytesseract OCR Extraction from {source_info}:\n{ocr_text.strip()}")
        else:
            logging.warning(f"pytesseract OCR extraction from {source_info} returned invalid text.")

        # Groq llama-3.2-11b-vision-preview extraction
        groq_text = await asyncio.to_thread(extract_data_with_groq, processed_image_bytes, GROQ_API_KEY)
        if groq_text and is_text_valid(groq_text):
            extracted_data.append(f"Groq llama-3.2-11b-vision-preview Extraction from {source_info}:\n{groq_text.strip()}")
        elif groq_text:
            logging.warning(f"Groq llama-3.2-11b-vision-preview extraction from {source_info} returned invalid text.")

        # Combine extracted data
        if extracted_data:
            combined_text = "\n\n".join(extracted_data)
            logging.info(f"Successfully extracted and combined text from {source_info}.")
            return combined_text
        else:
            logging.warning(f"No valid data extracted from {source_info}.")
            return ""
    except Exception as e:
        logging.error(f"Failed to process image from {source_info}: {e}")
        return ""

# ------------------------ Audio Processing Functions ------------------------

SAMPLE_RATE = 24000

def float32_to_pcm16(audio_data):
    """Convert float32 numpy array to base64-encoded PCM16."""
    logging.debug("Converting microphone input to PCM16 format.")

    # Normalize and convert to PCM16
    pcm_audio = (audio_data * 32767).astype(np.int16)

    # Ensure audio data is contiguous
    pcm_audio = np.ascontiguousarray(pcm_audio)

    # Base64 encode the PCM16 data
    pcm_base64 = base64.b64encode(pcm_audio.tobytes()).decode('utf-8')

    return pcm_base64

def play_audio_from_openai(base64_audio):
    """Play audio received from OpenAI (PCM16 in base64 format)."""
    logging.info("Playing audio received from OpenAI.")

    try:
        # Decode base64 audio data
        decoded_audio = base64.b64decode(base64_audio)

        # Create an AudioSegment from raw PCM data
        audio_segment = AudioSegment.from_raw(
            BytesIO(decoded_audio),
            sample_width=2,  # PCM16 has a sample width of 2 bytes
            frame_rate=SAMPLE_RATE,
            channels=1
        )

        logging.debug(f"Audio segment length (ms): {len(audio_segment)}")

        # Extract raw audio data to play it
        raw_audio = np.array(audio_segment.get_array_of_samples())

        # Convert raw_audio from int16 to float32 between -1.0 and 1.0
        audio_float32 = raw_audio.astype(np.float32) / 32768.0

        # Play the audio on the default output device
        import sounddevice as sd
        sd.play(audio_float32, samplerate=SAMPLE_RATE)
        sd.wait()
    except Exception as e:
        logging.error(f"Error while playing audio: {e}")

# Helper function to convert base64 PCM16 back to float32 numpy array
def base64_to_float32(base64_audio):
    """Convert base64-encoded PCM16 audio to float32 numpy array."""
    try:
        # Decode base64 audio data
        decoded_audio = base64.b64decode(base64_audio)

        # Convert bytes to numpy int16 array
        pcm_audio = np.frombuffer(decoded_audio, dtype=np.int16)

        # Convert int16 to float32 between -1.0 and 1.0
        audio_float32 = pcm_audio.astype(np.float32) / 32768.0

        return audio_float32
    except Exception as e:
        logging.error(f"Error converting base64 audio: {e}")
        return None

async def send_audio_chunk_to_openai(ws, base64_audio, sample_rate):
    message = {
        "type": "input_audio_buffer.append",
        "audio": base64_audio,
        "audio_format": {
            "encoding": "pcm16",
            "sample_rate": sample_rate,
            "channels": 1
        }
    }
    await ws.send(json.dumps(message))
    logging.debug(f"Audio chunk sent to OpenAI. Length: {len(base64_audio)} bytes, Sample rate: {sample_rate}")
    
    
# ------------------------ OPENAI REALTIME API INTERACTION ------------------------

async def setup_session_openai(ws):
    """Configure the session with 'alloy' voice and function definitions."""
    logging.info("Setting up session with 'alloy' voice.")
    session_setup = {
        "type": "session.update",
        "session": {
            "voice": "alloy",
            "instructions": (
                "Your knowledge cutoff is 2023-10. You are a helpful AI assistant. "
                "You have access to a knowledge base and retrieve relevant information using the 'get_relevant_information' function. "
                "You must **always use** the retrieved information from the 'information' field of the function's output to provide detailed and helpful answers to the user's question. "
                "**Assume that the retrieved information is accurate and up-to-date, even if it contradicts your knowledge cutoff.** "
                "If the information is insufficient, inform the user politely."
            ),
            "tools": [
                {
                    "name": "get_relevant_information",
                    "type": "function",  # Add this line
                    "description": "Retrieve relevant information from the knowledge base based on the user's query.",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "query": {
                                "type": "string",
                                "description": "The user's question."
                            }
                        },
                        "required": ["query"]
                    }
                }
            ]
        }
    }
    await ws.send(json.dumps(session_setup))
    logging.debug("Session setup sent.")

async def connect_to_openai():
    """Establish connection and set the session configuration."""
    logging.info("Connecting to OpenAI Realtime API.")
    headers = {
        "Authorization": f"Bearer {OPENAI_API_KEY}",
        "OpenAI-Beta": "realtime=v1"
    }
    try:
        ws = await websockets.connect(REALTIME_API_URL, extra_headers=headers, compression=None)
        await setup_session_openai(ws)  # Configure the session with Alloy voice and function definitions
        logging.info("Successfully connected to OpenAI.")
        return ws
    except Exception as e:
        logging.error("Failed to connect to OpenAI: %s", str(e))
        raise

# Add a global variable to track if a response is active
response_active_openai = False

async def receive_and_forward_responses_openai(ws, websocket):
    """Receive responses from OpenAI and forward them to the client."""
    global response_active_openai
    function_calls_in_progress = {}
    try:
        async for message in ws:
            try:
                response = json.loads(message)
                response_type = response.get('type')
                logging.debug(f"Received response of type: {response_type}")

                if response_type == 'error':
                    error_info = response.get('error', {})
                    logging.error(f"Error from OpenAI: {error_info}")
                    await websocket.send_text(json.dumps({
                        "type": "error",
                        "message": error_info.get('message', 'Unknown error')
                    }))
                    response_active_openai = False

                elif response_type == 'conversation.item.created':
                    item = response.get('item', {})
                    if item.get('type') == 'function_call':
                        function_name = item.get('name')
                        call_id = item.get('call_id')
                        function_calls_in_progress[call_id] = {
                            'name': function_name,
                            'arguments': ''
                        }

                elif response_type == 'response.function_call_arguments.delta':
                    call_id = response.get('call_id')
                    delta = response.get('delta', '')
                    if call_id in function_calls_in_progress:
                        function_calls_in_progress[call_id]['arguments'] += delta

                elif response_type == 'response.function_call_arguments.done':
                    call_id = response.get('call_id')
                    if call_id in function_calls_in_progress:
                        function_name = function_calls_in_progress[call_id]['name']
                        function_arguments = function_calls_in_progress[call_id]['arguments']
                        if function_name == 'get_relevant_information':
                            arguments = json.loads(function_arguments)
                            query = arguments.get('query', '')
                            relevant_paragraphs = get_relevant_paragraphs(query)
                            paragraphs_text = '\n\n'.join(relevant_paragraphs)
                            function_output = {"information": paragraphs_text}
                            await ws.send(json.dumps({
                                "type": "conversation.item.create",
                                "item": {
                                    "type": "function_call_output",
                                    "call_id": call_id,
                                    "output": json.dumps(function_output)
                                }
                            }))
                            await ws.send(json.dumps({"type": "response.create", "response": {"modalities": ["audio", "text"]}}))
                            response_active_openai = True
                            del function_calls_in_progress[call_id]

                elif response_type == 'response.audio.delta':
                    delta = response.get('delta')
                    if isinstance(delta, str):
                        await websocket.send_text(json.dumps({"type": "audio_response", "audio": delta}))

                elif response_type == 'response.text.delta':
                    delta = response.get('delta', '')
                    if isinstance(delta, str):
                        await websocket.send_text(json.dumps({"type": "text_response", "text": delta}))

                elif response_type == 'response.done':
                    logging.info("Response processing completed.")
                    response_active_openai = False

            except json.JSONDecodeError as e:
                logging.error(f"JSON decode error: {e}")

    except Exception as e:
        logging.error(f"Unhandled exception in OpenAI response handling: {e}")

# ------------------------ RAG FUNCTIONS ------------------------

def get_relevant_paragraphs(query: str) -> List[str]:
    """Retrieves relevant paragraphs from the vector index based on the query."""
    logging.debug("Searching for relevant documents using llama_index.")
    retriever = vector_index.as_retriever(similarity_top_k=3)
    retrieved_nodes = retriever.retrieve(query)  # Pass query string
    relevant_paragraphs = [node.node.get_content() for node in retrieved_nodes]
    logging.info(f"Retrieved {len(relevant_paragraphs)} relevant paragraphs.")
    return relevant_paragraphs

# ------------------------ Completion Generation Function ------------------------

def get_llama_completion(prompt: str) -> str:
    """
    Generates a text completion using NVIDIA's Llama model via NVIDIA's API.

    Parameters:
        prompt (str): The prompt for the Llama model.

    Returns:
        str: The generated text completion.
    """
    try:
        # NVIDIA API endpoint for completions
        url = f"{LLAMA_URL}/chat/completions"

        # Set up the headers with the API key
        headers = {
            "Authorization": f"Bearer {NVIDIA_API_KEY_LLAMA}",
            "Content-Type": "application/json",
        }

        # Prepare the payload
        payload = {
            "model": "meta/llama-3.1-8b-instruct",
            "messages": [{"role": "user", "content": prompt}],
            "max_tokens": 1024,
            "temperature": 0.2,
            "top_p": 0.7,
            "stream": False,
        }

        # Make the POST request
        response = requests.post(url, headers=headers, json=payload)

        # Check for successful response
        if response.status_code != 200:
            logging.error(f"Failed to get completion: {response.status_code} {response.text}")
            return ""

        response_json = response.json()
        completion = response_json["choices"][0]["message"]["content"].strip()
        logging.info(f"Generated completion for prompt: {prompt[:50]}... with length {len(completion)}")
        return completion
    except Exception as e:
        logging.error(f"Unexpected error during completion generation: {e}")
        return ""

# ------------------------ FastAPI Endpoints ------------------------

@app.get("/", response_class=HTMLResponse)
async def get_home(request: Request):
    """Serve the home page."""
    return templates.TemplateResponse("index.html", {"request": request})

@app.get("/contact", response_class=HTMLResponse)
async def get_contact(request: Request):
    """Serve the contact page."""
    return templates.TemplateResponse("contact.html", {"request": request})

@app.get("/processed_files", response_class=JSONResponse)
async def list_processed_files():
    """List all processed .txt files."""
    try:
        # List all .txt files in the DATABASE_DIR
        files = [f for f in os.listdir(DATABASE_DIR) if f.endswith('.txt')]

        if not files:
            # No processed files found
            return {"processed_files": [], "message": "No processed files found. Please upload a file to start processing."}

        return {"processed_files": files, "message": "Processed files listed successfully."}

    except Exception as e:
        logging.error(f"Error listing processed files: {e}")
        return JSONResponse(content={"error": "Error fetching processed files."}, status_code=500)

@app.post("/process_file", response_class=JSONResponse)
async def process_file_endpoint(file: UploadFile = File(...)):
    """Endpoint to upload and process a file."""
    file_content = await file.read()
    file_extension = file.filename.split('.')[-1].lower()

    try:
        # Extract text based on file type
        extracted_text = await process_document(file_content, file.content_type, file_extension, filename=file.filename)
        if not extracted_text or "Failed" in extracted_text or "Unsupported" in extracted_text:
            logging.warning(f"Failed processing file: {file.filename}. Unsupported file type or error in extraction.")
            return JSONResponse(content={"success": False, "error": "Failed to extract text from the file."}, status_code=400)

        # Save entire text as one .txt file with tool usage information
        doc_id = os.path.splitext(file.filename)[0]  # Remove extension for doc_id
        save_processed_document_to_db(doc_id, extracted_text)

        # Index the document
        index_document(extracted_text, doc_id=doc_id)

        logging.info(f"Extracted and indexed text from {file.filename} successfully.")
        return JSONResponse(content={"success": True, "message": f"File '{file.filename}' processed successfully!"}, status_code=200)

    except Exception as e:
        logging.error(f"Unexpected error processing file '{file.filename}': {e}")
        return JSONResponse(content={"success": False, "error": "Failed to process file due to an internal error."}, status_code=500)

async def process_document(file_content: bytes, content_type: str, file_extension: str, filename: str) -> str:
    """Process the uploaded document based on its file type."""
    if file_extension == "pdf":
        return await extract_text_from_pdf_with_ocr(file_content, filename)
    elif file_extension == "pptx":
        return await process_pptx(file_content, filename)
    elif file_extension == "txt":
        try:
            # Attempt to decode the txt file using UTF-8 encoding
            text = file_content.decode('utf-8')
            logging.info(f"Successfully decoded .txt file {filename}.")
            return f"Direct Text Extraction from {filename}:\n{text}"
        except UnicodeDecodeError:
            logging.error(f"Failed to decode .txt file {filename} with UTF-8 encoding.")
            try:
                # Fallback to ISO-8859-1 encoding
                text = file_content.decode('iso-8859-1')
                logging.info(f"Successfully decoded .txt file {filename} with ISO-8859-1 encoding.")
                return f"Direct Text Extraction from {filename} (ISO-8859-1):\n{text}"
            except Exception as e:
                logging.error(f"Failed to decode .txt file {filename} with ISO-8859-1 encoding: {e}")
                return f"Direct Text Extraction from {filename}:\nFailed to decode text."
    elif file_extension == "ppt":
        logging.warning(f"`.ppt` files are not supported for {filename}. Please convert to `.pptx` format.")
        return f"Unsupported file type: {file_extension.upper()} for {filename}."
    elif content_type.startswith("image/") or file_extension.lower() in ["jpg", "jpeg", "png", "gif", "bmp", "tiff", "webp"]:
        return await convert_and_extract_text_from_image(file_content, filename)
    elif file_extension in ["rtf", "doc", "docx", "odt"]:
        return extract_text_with_textract(file_content, file_extension, filename)
    elif file_extension.lower() in ["mp3", "mp4", "mpeg", "mpga", "m4a", "wav", "webm"]:
        return await process_audio_file(file_content, filename)
    else:
        logging.warning(f"Unsupported file type: {file_extension} for {filename}.")
        return f"Unsupported file type: {file_extension} for {filename}."

async def extract_text_from_pdf_with_ocr(pdf_content: bytes, filename: str) -> str:
    """Extract text from a PDF using pytesseract OCR and Groq's llama-3.2-11b-vision-preview."""
    try:
        images = convert_from_bytes(
            pdf_content, dpi=300, fmt='png'
        )
        text_content = []
        for page_number, img in enumerate(images, start=1):
            with BytesIO() as img_byte_io:
                img.save(img_byte_io, format='PNG', dpi=(300, 300))
                image_bytes = img_byte_io.getvalue()
            page_text = await process_image_content(image_bytes, source_info=f"{filename} - Page {page_number}")
            if is_text_valid(page_text):
                text_content.append(page_text)
        full_text = "\n\n".join(text_content)
        logging.info(f"Extracted text from PDF file {filename} using pytesseract OCR and Groq's llama-3.2-11b-vision-preview.")
        return full_text
    except Exception as e:
        logging.error(f"Failed to extract text from PDF {filename}: {e}")
        return ""

async def process_pptx(file_content: bytes, filename: str) -> str:
    """Process a PPTX file and extract text, including processing images."""
    try:
        prs = Presentation(BytesIO(file_content))
        text_content = []
        for slide_number, slide in enumerate(prs.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(f"Text Extraction from {filename} - Slide {slide_number}:\n{shape.text}")
                if shape.has_table:
                    table_data = [[cell.text for cell in row.cells] for row in shape.table.rows]
                    slide_text.append(f"Table Extraction from {filename} - Slide {slide_number}:\n{json.dumps(table_data, indent=2)}")
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image_bytes = shape.image.blob
                        image_text = await process_image_content(image_bytes, source_info=f"{filename} - Slide {slide_number} Image")
                        if is_text_valid(image_text):
                            slide_text.append(f"Image Extraction from {filename} - Slide {slide_number} Image:\n{image_text}")
                        else:
                            slide_text.append(f"Image Extraction from {filename} - Slide {slide_number} Image:\nNo valid text extracted from image.")
                    except Exception as e:
                        logging.error(f"Failed to process image in slide {slide_number}: {e}")
                        slide_text.append(f"Image Extraction from {filename} - Slide {slide_number} Image:\nError processing embedded image.")
            if slide_text:
                slide_content = "\n\n".join(slide_text)
                text_content.append(slide_content)
            else:
                logging.warning(f"No content extracted from slide {slide_number}.")
                text_content.append(f"Slide {slide_number} Content Extraction from {filename}:\nNo content extracted.")
        full_text = "\n\n".join(text_content)
        logging.info(f"Extracted text from PPTX file {filename} with detailed content.")
        return full_text
    except Exception as e:
        logging.error(f"Failed to extract text from PPTX file {filename}: {e}")
        return ""

async def convert_and_extract_text_from_image(image_content: bytes, filename: str) -> str:
    """Convert image and extract text using pytesseract OCR and Groq's llama-3.2-11b-vision-preview."""
    try:
        text = await process_image_content(image_content, filename)
        return text
    except Exception as e:
        logging.error(f"Failed to extract text from image {filename}: {e}")
        return ""

def extract_text_with_textract(file_content: bytes, file_extension: str, filename: str) -> str:
    """Extract text from various document formats using textract with tool usage information."""
    try:
        # Check if the file is RTF
        if file_extension.lower() == "rtf":
            # Use striprtf for RTF files
            text = rtf_to_text(file_content.decode('utf-8', errors='ignore'))
            logging.info(f"Extracted text from RTF file {filename} using striprtf.")
            return f"striprtf Extraction from {filename}:\n{text}"
        else:
            if file_extension.lower() == "odt":
                # Save to a temporary file
                with tempfile.NamedTemporaryFile(delete=False, suffix='.odt') as tmp:
                    tmp.write(file_content)
                    temp_path = tmp.name
                try:
                    text = textract.process(temp_path, extension=file_extension).decode("utf-8")
                finally:
                    os.unlink(temp_path)  # Ensure the temp file is deleted
            else:
                # Use textract for other supported file types
                text = textract.process(BytesIO(file_content), extension=file_extension).decode("utf-8")
            logging.info(f"Extracted text from {file_extension.upper()} file {filename} using textract.")
            return f"textract Extraction from {filename}:\n{text}"
    except Exception as e:
        logging.error(f"Failed to extract text from {file_extension.upper()} file {filename} using {'striprtf' if file_extension.lower() == 'rtf' else 'textract'}: {e}")
        return f"{'striprtf' if file_extension.lower() == 'rtf' else 'textract'} Extraction from {filename}:\nFailed to extract text."

async def process_audio_file(file_content: bytes, filename: str) -> str:
    """Processes an audio file using Groq's whisper-large-v3-turbo model.
    Returns the transcribed text."""
    try:
        # Since the Groq client is synchronous, use asyncio.to_thread to run it asynchronously
        def transcribe_audio():
            client = Groq(api_key=GROQ_API_KEY)
            transcription = client.audio.transcriptions.create(
                file=(filename, file_content),
                model="whisper-large-v3-turbo",
                response_format="json",
                # Optional parameters like 'prompt', 'language', 'temperature' can be added here
            )
            return transcription.text

        transcribed_text = await asyncio.to_thread(transcribe_audio)

        if is_text_valid(transcribed_text):
            logging.info(f"Successfully transcribed audio file {filename}.")
            return f"Groq whisper-large-v3-turbo Transcription from {filename}:\n{transcribed_text}"
        else:
            logging.warning(f"Transcribed text from {filename} is invalid.")
            return ""
    except Exception as e:
        logging.error(f"Failed to transcribe audio file {filename}: {e}")
        return ""
# ------------------------ WebSocket Endpoint for Q&A ------------------------

class QuestionMessage(BaseModel):
    type: str
    user_input: str

    @property
    def question(self):
        return self.user_input

@app.websocket("/llama_ws")
async def websocket_endpoint_llama(websocket: WebSocket):
    """WebSocket endpoint for real-time Q&A using similarity search."""
    await websocket.accept()
    logging.info("WebSocket connection established for /llama_ws.")
    is_closed = False

    try:
        while True:
            try:
                data = await websocket.receive_text()
                logging.debug(f"Received data: {data}")  # Log the entire message
            except WebSocketDisconnect:
                logging.info("WebSocket disconnected by client.")
                break
            except Exception as e:
                logging.error(f"Error receiving data from client: {e}")
                break

            # Validate the incoming message
            try:
                message = QuestionMessage.parse_raw(data)
                question = message.question
            except ValidationError as ve:
                logging.warning(f"Invalid message format: {ve}")
                await websocket.send_text(json.dumps({"error": "Invalid message format. 'type' and 'user_input' fields are required."}))
                continue  # Skip to the next iteration

            logging.info(f"Received question: {question}")

            if not vector_index:
                logging.error("VectorStoreIndex is not initialized. Cannot perform retrieval.")
                await websocket.send_text(json.dumps({"answer": "Internal error: Index not initialized."}))
                continue

            # Retrieve relevant paragraphs
            retriever = vector_index.as_retriever(similarity_top_k=5)
            retrieved_nodes = retriever.retrieve(question)

            if not retrieved_nodes:
                await websocket.send_text(json.dumps({"answer": "No relevant information found."}))
                continue

            # Prepare context from retrieved nodes
            context = "\n\n".join([node_with_score.node.get_content() for node_with_score in retrieved_nodes])

            # Prepare the messages for the model
            messages = [
                {"role": "system", "content": "Use the following information to answer the question."},
                {"role": "user", "content": f"Question: {question}\n\nContext: {context}"}
            ]

            # Prepare the payload for the NVIDIA API
            payload = {
                "messages": messages,
                "model": "meta/llama-3.1-8b-instruct",
                "max_tokens": 1024,
                "temperature": 0.2,
                "top_p": 0.7,
                "stream": True
            }

            try:
                headers = {
                    "Authorization": f"Bearer {NVIDIA_API_KEY_LLAMA.strip()}",
                    "Content-Type": "application/json",
                    "Accept": "application/json"
                }

                async with httpx.AsyncClient() as client:
                    async with client.stream("POST", f"{LLAMA_URL}/chat/completions", headers=headers, json=payload) as response:
                        if response.status_code == 401:
                            # Handle unauthorized error specifically
                            logging.error("Unauthorized: Check your API key.")
                            await websocket.send_text(json.dumps({"error": "Unauthorized access. Please check your API key."}))
                            break

                        response.raise_for_status()  # Raises error for other 4xx/5xx

                        async for line in response.aiter_lines():
                            line = line.strip()
                            if line == "[DONE]":
                                break  # Exit the loop when the DONE marker is reached

                            # Handle cases where the line contains data with "data: "
                            if line.startswith("data: "):
                                line = line[6:]

                            # Only attempt to parse lines that look like JSON
                            if line:
                                try:
                                    chunk = json.loads(line)
                                    content = chunk.get("choices", [{}])[0].get("delta", {}).get("content", "")
                                    if content:
                                        await websocket.send_text(json.dumps({"answer": content}))
                                        logging.info(f"Sent answer chunk: {content}")
                                except json.JSONDecodeError:
                                    logging.warning(f"Skipping non-JSON line: {line}")
            except httpx.HTTPStatusError as http_err:
                logging.error(f"HTTP error during stream: {http_err}")
                await websocket.send_text(json.dumps({"error": f"Failed to generate answer. Status Code: {http_err.response.status_code}"}))
            except Exception as e:
                logging.error(f"Error generating completion: {e}")
                await websocket.send_text(json.dumps({"error": "Failed to generate answer."}))

    except Exception as e:
        logging.error(f"Unhandled exception in WebSocket communication: {e}")
    finally:
        if not is_closed:
            try:
                await websocket.close()
                logging.info("WebSocket connection closed.")
            except Exception as e:
                logging.warning(f"Error closing WebSocket: {e}")
            finally:
                is_closed = True
                

# ------------------------ WebSocket Endpoint for OpenAI Realtime API ------------------------

@app.websocket("/openai_ws")
async def websocket_endpoint_openai(websocket: WebSocket):
    """WebSocket endpoint for OpenAI Realtime API interactions (Audio Only)."""
    global response_active_openai
    logging.info("WebSocket connection established for OpenAI Realtime API.")
    await websocket.accept()
    
    try:
        ws_openai = await connect_to_openai()
        
        # Create an asyncio task to receive and forward responses from OpenAI to the client
        receive_task = asyncio.create_task(receive_and_forward_responses_openai(ws_openai, websocket))
        
        # Trigger a response immediately to enable continuous interaction
        await ws_openai.send(json.dumps({"type": "response.create"}))

        while True:
            # Receive messages from the client
            try:
                data = await websocket.receive_text()
            except WebSocketDisconnect:
                logging.info("WebSocket disconnected by client.")
                break
            except Exception as e:
                logging.error(f"Error receiving data from client: {e}")
                break

            message = json.loads(data)
            message_type = message.get('type')
    
            if message_type == 'audio_chunk':
                # Receive audio data from the client
                base64_audio = message.get('audio')
                if base64_audio:
                    # Send the audio data to OpenAI
                    await send_audio_chunk_to_openai(ws_openai, base64_audio)
                else:
                    logging.error("No audio data received from client.")
                    error_message = {
                        "type": "error",
                        "message": "No audio data received."
                    }
                    await websocket.send_text(json.dumps(error_message))
            elif message_type == 'audio_stop':
                # The "audio_stop" is redundant in VAD mode and can be removed
                logging.warning("Received audio_stop, which is unnecessary in server VAD mode.")
            else:
                logging.error(f"Unknown message type received: {message_type}")
                error_message = {
                    "type": "error",
                    "message": f"Unknown message type: {message_type}. Only 'audio_chunk' is supported in server VAD mode."
                }
                await websocket.send_text(json.dumps(error_message))
    
    except Exception as e:
        logging.error(f"Error during WebSocket connection with OpenAI Realtime API: {e}")
    
    finally:
        # Cancel the background receive task with specific handling for CancelledError
        try:
            receive_task.cancel()
            await receive_task
        except asyncio.CancelledError:
            logging.info("Receive task was cancelled as expected during WebSocket cleanup.")
        except Exception as e:
            logging.warning(f"Error canceling receive_task: {e}")
    
        # Close the OpenAI WebSocket connection
        try:
            await ws_openai.close()
            logging.info("Closed OpenAI WebSocket connection.")
        except Exception as e:
            logging.warning(f"Error closing OpenAI WebSocket: {e}")
    
        # Close the client WebSocket connection if it's still open
        try:
            if websocket.client_state == "CONNECTED":
                await websocket.close()
                logging.info("Closed client WebSocket connection.")
        except RuntimeError as e:
            logging.warning(f"Attempted to close WebSocket but it was already closed: {e}")
        except Exception as e:
            logging.warning(f"Error closing client WebSocket: {e}")
    
        logging.info("WebSocket connection with OpenAI Realtime API closed.")


# ------------------------ send_audio_chunk_to_openai Function Update ------------------------

async def send_audio_chunk_to_openai(ws, base64_audio):
    # Only include "audio" in the message payload without "sample_rate" or "audio_format"
    message = {
        "type": "input_audio_buffer.append",
        "audio": base64_audio
    }
    
    await ws.send(json.dumps(message))
    logging.debug(f"Audio chunk sent to OpenAI. Length: {len(base64_audio)} bytes")
    
# ------------------------ receive_and_forward_responses_openai Update ------------------------

async def receive_and_forward_responses_openai(ws, websocket):
    """Receive responses from OpenAI and forward them to the client."""
    global response_active_openai
    function_calls_in_progress = {}
    try:
        async for message in ws:
            logging.debug(f"Received message from OpenAI: {message}")
            try:
                response = json.loads(message)
                response_type = response.get('type')
                logging.debug(f"Received response of type: {response_type}")

                if response_type == 'error':
                    error_info = response.get('error', {})
                    logging.error(f"Error from OpenAI: {error_info}")
                    await websocket.send_text(json.dumps({
                        "type": "error",
                        "message": error_info.get('message', 'Unknown error')
                    }))
                    response_active_openai = False

                elif response_type == 'conversation.item.created':
                    item = response.get('item', {})
                    if item.get('type') == 'function_call':
                        function_name = item.get('name')
                        call_id = item.get('call_id')
                        function_calls_in_progress[call_id] = {
                            'name': function_name,
                            'arguments': ''
                        }

                elif response_type == 'response.function_call_arguments.delta':
                    call_id = response.get('call_id')
                    delta = response.get('delta', '')
                    if call_id in function_calls_in_progress:
                        function_calls_in_progress[call_id]['arguments'] += delta

                elif response_type == 'response.function_call_arguments.done':
                    call_id = response.get('call_id')
                    if call_id in function_calls_in_progress:
                        function_name = function_calls_in_progress[call_id]['name']
                        function_arguments = function_calls_in_progress[call_id]['arguments']
                        if function_name == 'get_relevant_information':
                            arguments = json.loads(function_arguments)
                            query = arguments.get('query', '')
                            relevant_paragraphs = get_relevant_paragraphs(query)
                            paragraphs_text = '\n\n'.join(relevant_paragraphs)
                            function_output = {"information": paragraphs_text}
                            await ws.send(json.dumps({
                                "type": "conversation.item.create",
                                "item": {
                                    "type": "function_call_output",
                                    "call_id": call_id,
                                    "output": json.dumps(function_output)
                                }
                            }))
                            await ws.send(json.dumps({"type": "response.create", "response": {"modalities": ["audio", "text"]}}))
                            response_active_openai = True
                            del function_calls_in_progress[call_id]

                elif response_type == 'response.audio.delta':
                    delta = response.get('delta')
                    if isinstance(delta, str):
                        await websocket.send_text(json.dumps({"type": "audio_response", "audio": delta}))

                elif response_type == 'response.text.delta':
                    delta = response.get('delta', '')
                    if isinstance(delta, str):
                        await websocket.send_text(json.dumps({"type": "text_response", "text": delta}))

                elif response_type == 'response.done':
                    logging.info("Response processing completed.")
                    response_active_openai = False

            except json.JSONDecodeError as e:
                logging.error(f"JSON decode error: {e}")

    except Exception as e:
        logging.error(f"Unhandled exception in OpenAI response handling: {e}")
        
# ------------------------ Main Entry Point ------------------------

if __name__ == "__main__":
    logging.info("Starting FastAPI app.")
    import uvicorn
    uvicorn.run(app, host="127.0.0.1", port=8000)